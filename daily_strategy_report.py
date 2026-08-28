"""
台指策略日報自動生成 v2.9
================================================================
每天早上 8:00 由 GitHub Actions 自動執行

輸出：./reports/台股策略日報_YYYY.MM.DD.pptx

安裝套件：
    pip install yfinance python-pptx pandas numpy matplotlib requests holidays
================================================================
v2.7 變更：
  ① 時區修正：所有 datetime 統一使用台灣時間 (UTC+8)
  ② 顏色修正：全面改為台股慣例（紅多綠空）
  ③ 雲端字型：自動載入 NotoSans 中文字型
  ④ Yahoo Finance end 參數加 1 天 buffer，確保抓到最新資料
================================================================
v2.8 變更（修正「同一段歷史、不同執行日產出不同交易明細」問題）：
  ① 行情快照凍結：./cache/market_data_cache.csv
     每日抓到的行情存檔；超過 DATA_FREEZE_TDAYS 個交易日的舊資料一律
     以快取為準，不再被 yfinance auto_adjust 的回溯調整（除權息還原）
     或 FinMind 事後校正改寫。
  ② 訊號分數凍結：./cache/signal_cache.csv（v2.8 版本，已於 v2.9 修正邏輯）
  ③ 所有快取邏輯均有 try/except 保護，讀寫失敗自動退回即時計算，
     不影響日報產出；首次執行（無快取）亦可正常運作。
  ④ 環境變數 REBUILD_CACHE=1 可強制清空快取重建。
================================================================
v2.9 變更（修正 v2.8「等成熟才鎖定」的邏輯錯誤）：
  ★ 核心修正：訊號分數不再「等滿 25 個交易日才凍結」，而是「第一次被
    算出來的當下就永久鎖定」（apply_signal_cache()）。原因：交易決策
    一旦在當天被實際算出、實際會拿去操作，就必須永久保留，不能因為
    之後資料補齊、模型看法改變，就回頭假裝那天做了不同的決定——
    這才符合「依照當時實際發出的訊號記錄交易」的要求。
  ① export_trades_excel() 移除「訊號狀態」篩選/標記，恢復顯示全部交易——
     因為現在每一筆都是鎖定後的最終結果，不再有「暫定」這個概念。
  ② 新增 bootstrap_signal_cache.py（獨立腳本，需另外執行一次）：
     首次部署 v2.9 時，signal_cache.csv 是空的，若不處理，第一次
     執行會用「今天」的完整資料回頭鎖定最近一段時間的分數，鎖到的
     會是「事後回顧版」而非「當時真正看得到的版本」。這支腳本用
     逐日往回截斷資料、重新計算的方式，把最近一段時間的分數還原成
     真正的「時間點決策」版本，跑完後就不需要再執行第二次。
  ③ generate_daily_report() 回傳值維持 (pptx_path, excel_path, state) 不變，
    run_and_send.py 無須修改。
================================================================
"""

import os, sys, warnings
warnings.filterwarnings("ignore")
from datetime import datetime, timedelta, timezone
from pathlib import Path

import numpy as np
import pandas as pd
import requests
import matplotlib
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import matplotlib.dates as mdates
from matplotlib.lines import Line2D

try:
    import yfinance as yf
    YF_OK = True
except ImportError:
    YF_OK = False
    print("⚠️  yfinance 未安裝：pip install yfinance")

# ── 台灣時區 ──────────────────────────────────────────────────
TW_TZ = timezone(timedelta(hours=8))

# ── 字型設定（雲端 + 本機都能用）────────────────────────────────
import urllib.request
import matplotlib.font_manager as fm

font_path = Path("NotoSans.otf")
# 若字型檔不存在，自動從 Google Fonts 下載
if not font_path.exists():
    try:
        print("  正在下載中文字型 (NotoSans)...")
        font_url = "https://github.com/googlefonts/noto-cjk/raw/main/Sans/OTF/TraditionalChinese/NotoSansCJKtc-Regular.otf"
        urllib.request.urlretrieve(font_url, str(font_path))
    except Exception as e:
        print(f"  ⚠️ 字型下載失敗: {e}")

if font_path.exists():
    prop = fm.FontProperties(fname=str(font_path))
    fm.fontManager.addfont(str(font_path))
    matplotlib.rcParams['font.sans-serif'] = [prop.get_name()] + matplotlib.rcParams['font.sans-serif']
else:
    for _f in ['Microsoft JhengHei', 'PingFang TC', 'Heiti TC',
               'Noto Sans CJK TC', 'Noto Sans CJK JP', 'Arial Unicode MS']:
        try:
            matplotlib.rcParams['font.sans-serif'] = [_f] + matplotlib.rcParams['font.sans-serif']
            break
        except Exception:
            pass
matplotlib.rcParams['axes.unicode_minus'] = False


# ═══════════════════════════════════════════════════════════════
# Config
# ═══════════════════════════════════════════════════════════════
SCRIPT_DIR    = Path(os.getcwd())
TEMPLATE_PATH = SCRIPT_DIR / "daily_template.pptx"
OUTPUT_DIR    = SCRIPT_DIR / "reports"

FINMIND_TOKEN = "eyJ0eXAiOiJKV1QiLCJhbGciOiJIUzI1NiJ9.eyJ1c2VyX2lkIjoia3VvODYwMSIsImVtYWlsIjoic29sZGllcjg2MTAwQGdtYWlsLmNvbSIsInRva2VuX3ZlcnNpb24iOjB9._5JgdrkR3h3ogK7zaxW1t7R4UxB0rbR-_aZUm3z0HLQ"

LE, SE          = 3.5, 5.5
LX, SX          = 2.0, 2.0
EC_L, EC_S      = 5, 2
COST            = 0.0005
DYN_WINDOW      = 60
DYN_MIN_TRIG    = 3
DYN_LO, DYN_HI  = 0.5, 1.5
FWD_DAYS        = 20

# ── ★ v2.8 快照凍結設定 ────────────────────────────────────────
CACHE_DIR    = SCRIPT_DIR / "cache"
MARKET_CACHE = CACHE_DIR / "market_data_cache.csv"   # 行情原始資料快照
SIGNAL_CACHE = CACHE_DIR / "signal_cache.csv"        # ml / ms 訊號分數快照

# ★ v2.9：訊號分數不再等「成熟」才鎖定，而是第一次被算出來的當下就永久鎖定
#   （見 apply_signal_cache()）。這個常數保留下來純粹是給 bootstrap_signal_cache.py
#   當「回溯補建天數」的參考值：calc_dynamic 的 fwd20 需要 FWD_DAYS 個交易日
#   才算得完整，再加上 DYN_WINDOW 個交易日的視窗，理論上受影響的範圍大約是
#   FWD_DAYS + DYN_WINDOW ≈ 80 個交易日；這裡設寬鬆一點的參考值僅供說明，
#   實際回溯天數由 bootstrap 腳本的參數決定，跟這裡的日常凍結邏輯無關。
SIGNAL_MATURE_TDAYS = FWD_DAYS + DYN_WINDOW   # = 80（僅供參考，不再用於篩選）

# 行情凍結門檻：最近這幾個交易日的行情允許被新抓的資料更新（處理盤後校正），
# 更早的一律以快取為準，不再被回溯調整改寫。
DATA_FREEZE_TDAYS   = 10

BASE_L = {'fL1':1.0,'fL2':1.0,'fL3':2.0,'fL4':2.0,'fL5':2.0,'fL6':1.0,'fL7':0.5,
          'fL8':1.0,'fL9':1.0,'fL10':1.0,'fL11':0.5,'fL12':1.0}
BASE_S = {'fS1':1.5,'fS2':1.0,'fS3':2.0,'fS4':2.0,'fS5':2.0,'fS6':1.0,'fS7':0.5,
          'fS8':1.0,'fS9':1.0,'fS10':1.0,'fS11':0.5,'fS12':1.0,'fS13':1.0}


# ═══════════════════════════════════════════════════════════════
# 0. ★ v2.8 快照凍結工具
# ═══════════════════════════════════════════════════════════════
def reset_cache_if_requested():
    """環境變數 REBUILD_CACHE=1 時清空快取（用於策略邏輯變更後重建）"""
    if os.environ.get("REBUILD_CACHE", "").strip() not in ("1", "true", "True", "yes"):
        return
    for p in (MARKET_CACHE, SIGNAL_CACHE):
        try:
            if p.exists():
                p.unlink()
                print(f"  ♻️  REBUILD_CACHE=1：已刪除快取 {p.name}")
        except Exception as e:
            print(f"  ⚠️  刪除快取 {p.name} 失敗：{e}")


def _read_cache_csv(path: Path) -> pd.DataFrame:
    """讀取快取 CSV。float_precision="round_trip" 確保 float64 完全無損還原，
    否則微小的精度差會讓布林門檻翻轉，凍結就失去意義。舊版 pandas 不支援時自動退回。"""
    try:
        return pd.read_csv(path, float_precision="round_trip")
    except TypeError:
        return pd.read_csv(path)


def freeze_market_data(df: pd.DataFrame) -> pd.DataFrame:
    """
    ★ 行情快照凍結。
    - 讀取既有快取，將「已凍結區間」（快取最後日期往前推 DATA_FREEZE_TDAYS 之前）
      的數值覆蓋回本次抓到的資料，確保歷史行情永不被回溯調整改寫。
    - 快取有而本次沒抓到的日期/欄位，也會一併補回（防 YF 流量限制導致缺欄）。
    - 任何失敗都直接回傳原始 df，不影響日報產出。
    """
    out = df.copy()
    try:
        out["date"] = pd.to_datetime(out["date"])
        out = out.drop_duplicates(subset="date", keep="last").sort_values("date")
    except Exception as e:
        print(f"  ⚠️  行情快取：日期欄整理失敗，略過凍結（{e}）")
        return df

    try:
        CACHE_DIR.mkdir(parents=True, exist_ok=True)
    except Exception as e:
        print(f"  ⚠️  無法建立快取目錄，略過凍結（{e}）")
        return out.reset_index(drop=True)

    if MARKET_CACHE.exists():
        try:
            old = _read_cache_csv(MARKET_CACHE)
            old["date"] = pd.to_datetime(old["date"])
            old = old.drop_duplicates(subset="date", keep="last").sort_values("date")

            n_freeze = max(0, len(old) - DATA_FREEZE_TDAYS)
            frozen   = old.iloc[:n_freeze]

            if len(frozen) > 0:
                # combine_first：frozen 的值優先，缺值才用本次抓到的 out
                merged = frozen.set_index("date").combine_first(out.set_index("date"))
                merged = merged.sort_index()
                n_over = len(set(frozen["date"]) & set(out["date"]))
                print(f"  🔒 行情快取：{len(frozen)} 日已凍結"
                      f"（本次重疊 {n_over} 日採用快取值），"
                      f"最近 {DATA_FREEZE_TDAYS} 交易日仍即時更新")
                out = merged.reset_index()
            else:
                print(f"  ℹ️  行情快取尚未累積足夠天數（{len(old)} 日），本次全部即時")
        except Exception as e:
            print(f"  ⚠️  行情快取讀取失敗，本次改用即時資料（{e}）")

    try:
        # float_format="%.17g" 確保 float64 寫入/讀回完全無損（否則精度損失會讓閾值翻轉）
        out.to_csv(MARKET_CACHE, index=False, float_format="%.17g")
        print(f"  💾 行情快取已更新：{MARKET_CACHE.name}（{len(out)} 日）")
    except Exception as e:
        print(f"  ⚠️  行情快取寫入失敗（不影響本次日報）：{e}")

    return out.reset_index(drop=True)


def apply_signal_cache(d: pd.DataFrame, ml, ms):
    """
    ★ v2.9 訊號分數「當下即鎖定」（不是等成熟才鎖定）。

    背景：calc_dynamic() 用 fwd20 = shift(-FWD_DAYS)，這代表任何一天的分數，
    只要往回看的 60 天視窗裡有觸發樣本落在最近 ~FWD_DAYS+DYN_WINDOW 天內，
    分數就會隨資料越補越多而改變——這是模型設計本身的特性，不是 bug，
    永遠不會有「等到第 N 天之後就不會再變」這種乾淨的分界。

    但對「這一天實際做了什麼決策」這件事來說，重點從來不是這個分數準不準、
    成不成熟，而是「這是不是這一天第一次被算出來、當下實際會被拿去用的值」。
    一旦某一天已經被算過（不管是今天剛算出來、還是很久以前算出來的），
    這個值就要永久鎖住，之後任何一次重跑都不准再覆蓋——因為現實中的交易
    決策一旦做了就是做了，不能因為之後資料變多、模型「看法改變」就回頭
    假裝那天做了別的決定。

    做法：
      - 快取裡已經有的日期 → 一律套用快取值（無條件，不分成不成熟）。
      - 快取裡沒有的日期（代表這是第一次被算到，通常就是「今天」）
        → 用這次剛算出來的即時值，並立刻寫入快取、永久鎖定。

    回傳 (ml, ms)，任何失敗都回傳原值（退回沒有凍結的行為，不影響日報產出）。
    """
    ml = np.asarray(ml, dtype=float).copy()
    ms = np.asarray(ms, dtype=float).copy()

    try:
        dates = pd.to_datetime(d["date"]).reset_index(drop=True)
    except Exception as e:
        print(f"  ⚠️  訊號快取：日期欄讀取失敗，略過凍結（{e}）")
        return ml, ms

    try:
        CACHE_DIR.mkdir(parents=True, exist_ok=True)
    except Exception as e:
        print(f"  ⚠️  無法建立快取目錄，略過訊號凍結（{e}）")
        return ml, ms

    old = None
    if SIGNAL_CACHE.exists():
        try:
            old = _read_cache_csv(SIGNAL_CACHE)
            old["date"] = pd.to_datetime(old["date"])
            # keep="first"：最早寫入的（最早鎖定的）永遠優先，之後重複一律忽略
            old = old.drop_duplicates(subset="date", keep="first")

            mp_ml = dates.map(old.set_index("date")["ml"]).to_numpy(dtype=float)
            mp_ms = dates.map(old.set_index("date")["ms"]).to_numpy(dtype=float)
            hit   = (~np.isnan(mp_ml)) & (~np.isnan(mp_ms))

            n_diff = 0
            if hit.any():
                n_diff = int((np.abs(mp_ml[hit] - ml[hit]) > 1e-9).sum())
                ml[hit] = mp_ml[hit]
                ms[hit] = mp_ms[hit]

            print(f"  🔒 訊號快取：沿用 {int(hit.sum())} 日已鎖定分數"
                  f"（其中 {n_diff} 日與本次即時重算結果不同 → 已阻止歷史被改寫）")
        except Exception as e:
            print(f"  ⚠️  訊號快取讀取失敗，本次全部即時計算（{e}）")
            old = None

    # ★ 把「這次剛出現、快取裡還沒有」的日期立刻寫入快取鎖定——
    #   不再等成熟，第一次算出來的當下就是永久值。
    try:
        fresh = pd.DataFrame({"date": dates.values, "ml": ml, "ms": ms})
        if old is not None and len(old) > 0:
            combined = pd.concat([old[["date", "ml", "ms"]], fresh], ignore_index=True)
        else:
            combined = fresh
        # keep="first"：已鎖定的日期（在 old 裡、排在前面）永遠優先於這次剛算出來的重複值
        combined = combined.drop_duplicates(subset="date", keep="first").sort_values("date")
        n_new = len(combined) - (len(old) if old is not None else 0)
        combined.to_csv(SIGNAL_CACHE, index=False, float_format="%.17g")
        print(f"  💾 訊號快取已更新：{SIGNAL_CACHE.name}"
              f"（總計鎖定 {len(combined)} 日，本次新鎖定 {n_new} 日）")
    except Exception as e:
        print(f"  ⚠️  訊號快取寫入失敗（不影響本次日報）：{e}")

    return ml, ms


# ═══════════════════════════════════════════════════════════════
# 1. 資料抓取（Yahoo Finance + CSV fallback）
# ═══════════════════════════════════════════════════════════════
def fetch_market_data() -> pd.DataFrame:
    """Yahoo Finance；失敗時 fallback 本地 CSV"""
    if YF_OK:
        now_tw = datetime.now(TW_TZ).replace(tzinfo=None)
        specs = {
            "TWII"    : "^TWII",
            "SOX"     : "^SOX",
            "TSMC_TW" : "2330.TW",
            "TSM_US"  : "TSM",
            "ELEC"    : "0053.TW",
            "FIN"     : "0055.TW",
            "USDTWD"  : "USDTWD=X",
        }
        frames = {}

        # ★ 改用 yf.download() 批次下載，比逐一 Ticker.history() 穩定許多
        ticker_list = list(specs.values())
        ticker_names = list(specs.keys())
        try:
            import time
            print(f"  正在從 Yahoo Finance 下載資料...")
            # ★ 固定起始日 2019-01-01，原因：
            #   1. 策略成立日為 2021/06/16，MA120 需要 ~6 個月暖機期
            #   2. 若從 2021/06/01 抓，2021/06/16 時 MA60/MA120 尚未算滿，指標不準
            #   3. 從 2019 年起抓，到 2021/06/16 已有約 2.5 年暖機，指標完全穩定
            #   4. 績效統計仍從 INCEPTION_DATE=2021/06/16 開始，不受影響
            DATA_START = "2019-01-01"
            raw_all = yf.download(
                tickers=ticker_list,
                start=DATA_START,
                auto_adjust=True,
                progress=False,
                group_by="ticker",
            )
            if raw_all.empty:
                time.sleep(3)
                raw_all = yf.download(
                    tickers=ticker_list, start=DATA_START,
                    auto_adjust=True, progress=False, group_by="ticker"
                )
        except Exception as e:
            print(f"  ⚠️  yf.download 失敗: {e}")
            raw_all = pd.DataFrame()

        if not raw_all.empty:
            for name, ticker in specs.items():
                try:
                    if len(ticker_list) == 1:
                        tk_df = raw_all
                    else:
                        tk_df = raw_all[ticker] if ticker in raw_all.columns.get_level_values(0) else pd.DataFrame()
                    if tk_df.empty: continue
                    if tk_df.index.tz is not None:
                        tk_df.index = tk_df.index.tz_localize(None)
                    tk_df.index = pd.to_datetime(tk_df.index).normalize()
                    frames[name] = tk_df[["Close"]].rename(columns={"Close": name})
                    if name == "TSMC_TW" and "Volume" in tk_df.columns:
                        frames["TSMC_Vol"] = tk_df[["Volume"]].rename(columns={"Volume":"TSMC_Vol"})
                    if name == "TWII":
                        if "Open"   in tk_df.columns: frames["TWII_Open"] = tk_df[["Open"]].rename(columns={"Open":"TWII_Open"})
                        if "High"   in tk_df.columns: frames["TWII_High"] = tk_df[["High"]].rename(columns={"High":"TWII_High"})
                        if "Low"    in tk_df.columns: frames["TWII_Low"]  = tk_df[["Low"]].rename(columns={"Low":"TWII_Low"})
                except Exception as e:
                    print(f"  ⚠️  解析 {ticker}: {e}")

        # ★ 新增：YF 的 ^TWII 經常延遲一天或完全抓不到。
        # 這裡改呼叫 FinMind 官方 API 取得最新加權指數來補強，甚至在 YF 失敗時直接取代！
        try:
            start_str = (now_tw - timedelta(days=20)).strftime("%Y-%m-%d")
            url = "https://api.finmindtrade.com/api/v4/data"
            params = {"dataset":"TaiwanStockPrice", "data_id":"TAIEX", "start_date":start_str, "token":FINMIND_TOKEN}
            res = requests.get(url, params=params, headers={"User-Agent":"Mozilla/5.0"}, timeout=15)
            data = res.json()
            if data.get("msg") == "success" and data.get("data"):
                fm_df = pd.DataFrame(data["data"])
                fm_df["date"] = pd.to_datetime(fm_df["date"]).dt.normalize()
                fm_df = fm_df.set_index("date")

                # 若 YF 完全沒抓到，就用 FinMind 建立骨架
                if "TWII" not in frames: frames["TWII"] = pd.DataFrame(columns=["TWII"])
                if "TWII_Open" not in frames: frames["TWII_Open"] = pd.DataFrame(columns=["TWII_Open"])
                if "TWII_High" not in frames: frames["TWII_High"] = pd.DataFrame(columns=["TWII_High"])
                if "TWII_Low"  not in frames: frames["TWII_Low"]  = pd.DataFrame(columns=["TWII_Low"])

                for dt, row in fm_df.iterrows():
                    # FinMind 資料絕對準確，直接覆蓋 YF 的假資料或填補空缺
                    frames["TWII"].loc[dt, "TWII"] = float(row["close"])
                    frames["TWII_Open"].loc[dt, "TWII_Open"] = float(row["open"])
                    frames["TWII_High"].loc[dt, "TWII_High"] = float(row["max"])
                    frames["TWII_Low"].loc[dt, "TWII_Low"]  = float(row["min"])

                print(f"  ✓  FinMind 最新加權指數補強成功：{fm_df.index[-1].date()} -> {fm_df.iloc[-1]['close']:,.0f}")
        except Exception as e:
            print(f"  ⚠️  FinMind 補強失敗: {e}")

        # 合併前先去除任何可能的重複日期（YF 常見 bug 會導致 concat 崩潰）
        for name in list(frames.keys()):
            frames[name] = frames[name][~frames[name].index.duplicated(keep='last')]

        if frames and "TWII" in frames:
            df = pd.concat(frames.values(), axis=1).sort_index()

            # ★ 剔除 YF 產生的假 K 棒（當日未開盤前，YF 會生出一根開高低收完全相等的 K 棒）
            if all(c in df.columns for c in ["TWII", "TWII_Open", "TWII_High", "TWII_Low"]):
                fake_mask = (df["TWII"] == df["TWII_Open"]) & (df["TWII_High"] == df["TWII_Low"]) & (df["TWII"] == df["TWII_High"])
                df = df[~fake_mask]

            # ★ 關鍵修正：必須先剔除 TWII 沒有資料的日子，再 ffill()。
            # 否則如果昨天台股沒資料(或放假)而美股有開，會把「前天」的台股收盤價錯誤複製到「昨天」！
            df = df.dropna(subset=["TWII"])
            df = df.ffill()

            # USDTWD 合理性驗證（應介於 15~50 TWD/USD）
            if "USDTWD" in df.columns:
                rate = df["USDTWD"].iloc[-1]
                if not (15 < rate < 50):
                    print(f"  ⚠️  USDTWD 異常值 {rate:.4f}，自動取倒數修正")
                    df["USDTWD"] = 1.0 / df["USDTWD"]
                print(f"  ✓  匯率 USD/TWD = {df['USDTWD'].iloc[-1]:.2f}")

            # ADR 折溢價：1 TSM ADR = 5 股 2330.TW
            if all(c in df.columns for c in ["TSM_US","USDTWD","TSMC_TW"]):
                df["ADR_Premium"] = (df["TSM_US"] * df["USDTWD"] / 5.0
                                     / df["TSMC_TW"] - 1.0) * 100.0
            else:
                df["ADR_Premium"] = np.nan

            df = df.reset_index().rename(columns={"index":"date"})
            if "date" not in df.columns:
                df = df.rename(columns={df.columns[0]:"date"})

            # ★ 避免 Yahoo Finance 偷跑：早上盤前執行時，若撈到今天的假 K 棒必須剔除
            if now_tw.hour < 14:
                df = df[df["date"].dt.date < now_tw.date()].reset_index(drop=True)

            print(f"  ✓  Yahoo 最新資料日期：{df['date'].iloc[-1].date()}")
            print(f"  ✓  最終資料日期：{df['date'].iloc[-1].date()}")
            # 只要有 TWII 就可以繼續（其他欄位缺失時模型會自動降級）
            if "TWII" in df.columns:
                missing = [c for c in ["SOX","TSMC_TW","ELEC"] if c not in df.columns]
                if missing:
                    print(f"  ⚠️  部分欄位缺失（可能 YF 流量限制）: {missing}，計算時將忽略")
                return df

    csv_path = SCRIPT_DIR / "台指量化模型基礎資料表_包含電子金融.csv"
    if csv_path.exists():
        print(f"  ⚠️  Yahoo 失敗，改用 CSV: {csv_path.name}")
        return _load_from_csv(csv_path)
    raise RuntimeError("無法取得市場資料，可能是 Yahoo Finance 流量限制 (429)，請等待 30 分鐘後再試。")


def _load_from_csv(path) -> pd.DataFrame:
    raw = pd.read_csv(path)
    col_map = {
        '日期':'date','台灣加權指數開盤價':'TWII_Open','台灣加權指數收盤價':'TWII',
        '費半收盤價':'SOX','台積電收盤':'TSMC_TW','台積電成交量':'TSMC_Vol',
        '台積電ADR折溢價(%)':'ADR_Premium',
        '電子指數(0053)收盤價':'ELEC','金融指數(0055)收盤價':'FIN',
    }
    raw = raw.rename(columns=col_map)
    raw['date'] = pd.to_datetime(raw['date'])
    for c in raw.columns:
        if c != 'date' and raw[c].dtype == object:
            raw[c] = pd.to_numeric(
                raw[c].astype(str).str.replace(',','').str.replace('%',''), errors='coerce')
    raw = raw.dropna(subset=['TWII','TWII_Open']).ffill().reset_index(drop=True)
    fi_col = '外資合計買賣超金額(百萬)'
    if fi_col in raw.columns:
        raw['FI_Net_csv'] = pd.to_numeric(raw[fi_col], errors='coerce') / 100
    return raw


def fetch_foreign_investor(start_date: str) -> pd.Series:
    try:
        url    = "https://api.finmindtrade.com/api/v4/data"
        params = {"dataset":"TaiwanStockTotalInstitutionalInvestors",
                  "start_date":start_date,"token":FINMIND_TOKEN}
        res    = requests.get(url, params=params,
                              headers={"User-Agent":"Mozilla/5.0"}, timeout=15)
        data   = res.json()
        if data.get("msg") != "success": return pd.Series(dtype=float)
        dff    = pd.DataFrame(data["data"])
        mask   = dff["name"].str.contains("外資|Foreign_Investor", case=False, na=False) & \
                 ~dff["name"].str.contains("自營商|Dealer", case=False, na=False)
        f = dff[mask].copy()
        if f.empty: return pd.Series(dtype=float)
        f["Date"]   = pd.to_datetime(f["date"]).dt.normalize()
        f["FI_Net"] = (f["buy"].astype(float)-f["sell"].astype(float))/1e8
        return f.groupby("Date")["FI_Net"].sum()
    except Exception:
        return pd.Series(dtype=float)


# ═══════════════════════════════════════════════════════════════
# 2. 因子計算
# ═══════════════════════════════════════════════════════════════
def build_factors(df: pd.DataFrame, fi: pd.Series) -> pd.DataFrame:
    d = df.copy()
    d["MA5"]  = d["TWII"].rolling(5).mean()
    d["MA10"] = d["TWII"].rolling(10).mean()
    d["MA20"] = d["TWII"].rolling(20).mean()
    d["MA60"] = d["TWII"].rolling(60).mean()
    d["斜率"]  = (d["MA60"].diff(5) / d["MA60"].shift(5)) * 100
    d["乖離"]  = (d["TWII"] - d["MA60"]) / d["MA60"] * 100
    d["STD20"] = d["TWII"].rolling(20).std()
    d["BB上"]  = d["MA20"] + 2*d["STD20"]
    d["BB下"]  = d["MA20"] - 2*d["STD20"]
    δ = d["TWII"].diff()
    d["RSI"] = 100 - (100 / (1 +
        δ.clip(lower=0).ewm(com=13,adjust=False).mean() /
        (-δ.clip(upper=0)).ewm(com=13,adjust=False).mean().replace(0,np.nan)))
    d["MOM5"]    = (d["TWII"] / d["TWII"].shift(5) - 1) * 100
    d["SOX_MA20"]= d["SOX"].rolling(20).mean()
    d["SOX_MA60"]= d["SOX"].rolling(60).mean()
    d["TS_MA20"] = d["TSMC_TW"].rolling(20).mean()
    d["TS_VolMA"]= d["TSMC_Vol"].rolling(10).mean()
    d["EF"]      = d["ELEC"] / d["FIN"]
    d["EF_MA20"] = d["EF"].rolling(20).mean()
    d["EF_MA60"] = d["EF"].rolling(60).mean()
    if not fi.empty:
        d = d.set_index("date").join(fi.rename("FI_Net"), how="left").reset_index()
        d["FI_Net"]  = d["FI_Net"].ffill().fillna(0)
        d["FI_MA"]   = d["FI_Net"].rolling(120,min_periods=30).mean()
        d["FI_STD"]  = d["FI_Net"].rolling(120,min_periods=30).std().replace(0,np.nan)
        d["FI_Z"]    = (d["FI_Net"] - d["FI_MA"]) / d["FI_STD"]
        d["FI_5MA"]  = d["FI_Net"].rolling(5).mean()
    else:
        d["FI_Net"] = 0.0; d["FI_Z"] = 0.0; d["FI_5MA"] = 0.0
    d["ADR_MA"]  = d["ADR_Premium"].rolling(120,min_periods=30).mean()
    d["ADR_STD"] = d["ADR_Premium"].rolling(120,min_periods=30).std().replace(0,np.nan)
    d["ADR_Z"]   = (d["ADR_Premium"] - d["ADR_MA"]) / d["ADR_STD"]
    d["fL1"]  = ((d["TWII"]>d["MA60"]) & (d["斜率"]>0.1)).astype(float)
    d["fL2"]  = (d["EF"]>d["EF_MA20"]).astype(float)
    d["fL3"]  = (d["FI_Z"]>1.2).astype(float)
    d["fL4"]  = ((d["SOX"]>d["SOX_MA20"]) & (d["SOX"]>d["SOX_MA60"])).astype(float)
    d["fL5"]  = (d["ADR_Z"]>0.8).astype(float)
    d["fL6"]  = (d["TSMC_TW"]>d["TS_MA20"]).astype(float)
    d["fL7"]  = (d["TSMC_Vol"]>1.5*d["TS_VolMA"]).astype(float)
    d["fL8"]  = (d["乖離"]<-8).astype(float)
    d["fL9"]  = (d["RSI"]<40).astype(float)
    d["fL10"] = (d["FI_5MA"]>0).astype(float)
    d["fL11"] = (d["TWII"]<d["BB下"]).astype(float)
    d["fL12"] = (d["MOM5"]>2).astype(float)
    d["fS1"]  = ((d["TWII"]<d["MA60"]) & (d["斜率"]<-0.1)).astype(float)
    d["fS2"]  = (d["EF"]<d["EF_MA20"]).astype(float)
    d["fS3"]  = (d["FI_Z"]<-1.2).astype(float)
    d["fS4"]  = ((d["SOX"]<d["SOX_MA20"]) & (d["SOX"]<d["SOX_MA60"])).astype(float)
    d["fS5"]  = (d["ADR_Z"]<-1.0).astype(float)
    d["fS6"]  = (d["TSMC_TW"]<d["TS_MA20"]).astype(float)
    d["fS7"]  = (d["TSMC_Vol"]>1.5*d["TS_VolMA"]).astype(float)
    d["fS8"]  = (d["乖離"]>8).astype(float)
    d["fS9"]  = (d["RSI"]>55).astype(float)
    d["fS10"] = (d["FI_5MA"]<0).astype(float)
    d["fS11"] = (d["TWII"]>d["BB上"]).astype(float)
    d["fS12"] = (d["TWII"]<d["MA60"]).astype(float)
    d["fS13"] = (d["MOM5"]<-2).astype(float)
    d["fwd20"] = d["TWII"].shift(-FWD_DAYS) / d["TWII"] - 1
    return d


def calc_dynamic(factor_arr, fwd_arr, is_long):
    N = len(factor_arr); hits = np.full(N, 0.5)
    for i in range(DYN_WINDOW, N):
        fw = factor_arr[i-DYN_WINDOW:i]; rw = fwd_arr[i-DYN_WINDOW:i]
        trig = fw == 1.0
        if trig.sum() >= DYN_MIN_TRIG:
            rs = rw[trig]; rs = rs[~np.isnan(rs)]
            if len(rs) >= DYN_MIN_TRIG:
                hits[i] = (rs>0).mean() if is_long else (rs<0).mean()
    return hits


def compute_scores(d):
    fwd   = d["fwd20"].values
    dyn_L = {k: calc_dynamic(d[k].values, fwd, True)  for k in BASE_L}
    dyn_S = {k: calc_dynamic(d[k].values, fwd, False) for k in BASE_S}
    N = len(d); ml = np.zeros(N); ms = np.zeros(N)
    for k,w in BASE_L.items():
        ml += d[k].values * w * np.clip(dyn_L[k]*2, DYN_LO, DYN_HI)
    for k,w in BASE_S.items():
        ms += d[k].values * w * np.clip(dyn_S[k]*2, DYN_LO, DYN_HI)
    gL = ((d["fL4"]==1)&((d["fL5"]==1)|(d["fL3"]==1))&(d["EF"]>d["EF_MA60"])).values
    gS = ((d["fS4"]==1)&((d["fS5"]==1)|(d["fS3"]==1))&(d["EF"]<d["EF_MA20"])).values
    return ml, ms, gL, gS


def backtest(d, ml, ms, gL, gS):
    N = len(d)
    close = d["TWII"].values; open_ = d["TWII_Open"].values
    ma10  = d["MA10"].values; ma60  = d["MA60"].values
    intra = np.where(open_>0, close/open_-1, 0)
    onight = np.zeros(N)
    onight[1:] = np.where(close[:-1]>0, open_[1:]/close[:-1]-1, 0)
    daily = np.zeros(N)
    daily[1:]  = np.where(close[:-1]>0, close[1:]/close[:-1]-1, 0)
    pos = np.zeros(N); cur=0.0; ec=0; quick=0
    for i in range(N):
        if cur == 0:
            if ml[i]>=LE and gL[i]:    cur=1.;  ec=0; quick=0
            elif ms[i]>=SE and gS[i]:  cur=-1.; ec=0; quick=0
        elif cur == 1:
            esig = ml[i]<LX or not gL[i] or close[i]<ma60[i]
            ec   = ec+1 if esig else 0
            if ec>=EC_L: cur=0.; ec=0
        else:
            if close[i]>ma10[i]: quick+=1
            else: quick=0
            if quick>=1:
                cur=0.; ec=0; quick=0; pos[i]=cur; continue
            esig = ms[i]<SX or not gS[i] or close[i]>ma60[i]
            ec   = ec+1 if esig else 0
            if ec>=EC_S: cur=0.; ec=0
        pos[i] = cur
    exp  = np.roll(pos,1); exp[0]=0
    expp = np.roll(exp,1); expp[0]=0
    ret  = np.zeros(N)
    me=(exp!=0)&(expp==0);           ret[me]=exp[me]*intra[me]
    mh=(exp!=0)&(expp==exp);         ret[mh]=exp[mh]*daily[mh]
    mx=(exp==0)&(expp!=0);           ret[mx]=expp[mx]*onight[mx]
    mr=(exp!=0)&(expp!=0)&(exp!=expp)
    ret[mr]=expp[mr]*onight[mr]+exp[mr]*intra[mr]
    ret -= np.abs(np.diff(exp, prepend=0))*COST
    trades=[]; it=False; tr=[]; cd=0; ei=0
    for i in range(N):
        e=exp[i]
        if not it and e!=0:
            it=True; tr=[ret[i]]; cd=int(e); ei=i
            # ★ 進場價：使用執行日（訊號觸發後的隔天）的開盤價，與報酬計算一致
            trades.append({"entry_idx":i,"entry_date":d["date"].iloc[i],
                           "dir_code":cd,"entry_price":open_[i]})
        elif it and e!=0:
            tr.append(ret[i])
        elif it and e==0:
            tr.append(ret[i])
            p=np.prod(1+np.array(tr))-1
            # ★ 出場價：使用出場日的開盤價，與報酬計算（onight）一致
            trades[-1].update({"exit_idx":i,"exit_date":d["date"].iloc[i],
                               "exit_price":open_[i],"pct_return":p*100,
                               "n_days":i-ei,"is_win":p>0})
            it=False; tr=[]
    return {"daily_ret":ret,"exp":exp,"pos":pos,"cum":np.cumprod(1+ret),"trades":trades}


# ═══════════════════════════════════════════════════════════════
# 3. 狀態判斷
# ═══════════════════════════════════════════════════════════════
def determine_state(d, bt, ml, ms, gL, gS):
    i=len(d)-1; pt=int(bt["pos"][i]); py=int(bt["pos"][i-1]) if i>=1 else 0
    rec={( 1, 1):"多單續抱",( 1, 0):"多單建倉",( 1,-1):"空單轉多",
         ( 0, 1):"多單出場",( 0, 0):"空手觀望",( 0,-1):"空單出場",
         (-1, 1):"多單轉空",(-1, 0):"空單建倉",(-1,-1):"空單續抱",
         }.get((pt,py),"空手觀望")
    if   rec in ("多單續抱","多單建倉","空單轉多"): bias="偏多"
    elif rec in ("空單續抱","空單建倉","多單轉空"): bias="偏空"
    else:                                            bias="震盪"
    return {"recommendation":rec,"bias":bias,
            "ml":float(ml[i]),"ms":float(ms[i]),
            "gL":bool(gL[i]),"gS":bool(gS[i]),
            "pos_today":pt,"pos_yesterday":py,
            "last_close":float(d["TWII"].iloc[i]),
            "last_date":d["date"].iloc[i].date()}


def _calc_period_stats(d, bt, cutoff):
    """計算指定起始日之後的績效統計（通用子函式）"""
    mask = (d["date"] >= cutoff).values
    if mask.sum() < 5:
        return {"error": "資料不足"}
    sub_ret = bt["daily_ret"][mask]
    sub_exp = bt["exp"][mask]
    sc  = np.cumprod(1 + sub_ret)
    tot = (sc[-1] - 1) * 100
    yrs = mask.sum() / 252
    trades_p = [t for t in bt["trades"] if "exit_date" in t and t["entry_date"] >= cutoff]
    n_t = len(trades_p)
    n_l = sum(1 for t in trades_p if t["dir_code"] ==  1)
    n_s = sum(1 for t in trades_p if t["dir_code"] == -1)
    w_l = sum(1 for t in trades_p if t["dir_code"] ==  1 and t["is_win"])
    w_s = sum(1 for t in trades_p if t["dir_code"] == -1 and t["is_win"])
    w_t = sum(1 for t in trades_p if t["is_win"])
    return {
        "period":       (d["date"][mask].iloc[0].date(), d["date"].iloc[-1].date()),
        "n_trades":     n_t, "n_long": n_l, "n_short": n_s,
        "wr_all":       (w_t / n_t * 100) if n_t else 0,
        "wr_long":      (w_l / n_l * 100) if n_l else 0,
        "wr_short":     (w_s / n_s * 100) if n_s else 0,
        "strat_total":  tot,
        # ★ 未滿一年不年化（短期高報酬放大會誤導），回傳 None 讓表格顯示「—（未滿一年）」
        "strat_ann":    ((1 + tot / 100) ** (1 / yrs) - 1) * 100 if yrs >= 1.0 else None,
        "strat_mdd":    ((sc / np.maximum.accumulate(sc) - 1) * 100).min(),
        "strat_vol":    sub_ret.std() * np.sqrt(252) * 100,
        "strat_sharpe": sub_ret.mean() * 252 / (sub_ret.std() * np.sqrt(252) + 1e-9),
        "in_mkt":       (sub_exp != 0).mean() * 100,
        "in_long":      (sub_exp == 1).mean() * 100,
        "in_short":     (sub_exp == -1).mean() * 100,
    }


def compute_perf_stats(d, bt):
    """計算「成立至今（2021/06/16起）」＋今年以來(YTD) 績效統計，回傳 dict with keys '3yr' and 'ytd'"""
    # 起始日固定為策略成立日 2021/06/16
    INCEPTION_DATE = pd.Timestamp(2021, 6, 16)
    cutoff_ytd = pd.Timestamp(d["date"].iloc[-1].year, 1, 1)
    return {
        "3yr": _calc_period_stats(d, bt, INCEPTION_DATE),
        "ytd": _calc_period_stats(d, bt, cutoff_ytd),
    }


def compute_1yr_stats(d, bt):
    """近一年統計（供 Excel 匯出及終端機摘要輸出使用）"""
    return _calc_period_stats(d, bt, d["date"].iloc[-1] - pd.Timedelta(days=365))



# ═══════════════════════════════════════════════════════════════
# 4. 圖檔生成（台股慣例：紅多綠空）
# ═══════════════════════════════════════════════════════════════
def add_stats_table_to_slide(slide, stats_3yr, stats_ytd):
    """
    在 PPT slide 上建立可編輯的績效統計表（近三年 + 今年以來YTD 對照）。
    取代原本的 matplotlib 圖片，讓內容可在 PowerPoint 中直接編輯。
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── 定義表格每一列 ────────────────────────────────────────
    # (標籤, 取值函式, 用於顏色判斷的 key, 顏色規則)
    # 顏色規則: neutral / pos_red(正=紅負=綠) / neg_green(負=綠,MDD用) / wr(>=60%=紅)

    def _ann_text(s):
        """年化報酬：不足一年顯示 '—（未滿一年）'，避免短期高報酬誤導"""
        if "error" in s:
            return "N/A"
        v = s.get("strat_ann", None)
        if v is None:
            return "— (未滿一年)"
        return f"{v:+.2f}%"

    ROWS = [
        ("回測期間",
         lambda s: f"{s['period'][0]:%Y/%m/%d}~\n{s['period'][1]:%Y/%m/%d}"
                   if "error" not in s else "資料不足",
         None, "neutral"),
        ("交易筆數",
         lambda s: f"{s.get('n_trades', 0)} 筆" if "error" not in s else "N/A",
         None, "neutral"),
        ("多 / 空筆數",
         lambda s: f"{s.get('n_long', 0)} / {s.get('n_short', 0)}"
                   if "error" not in s else "N/A",
         None, "neutral"),
        ("整體勝率",
         lambda s: f"{s.get('wr_all', 0):.1f}%"   if "error" not in s else "N/A",
         "wr_all",       "wr"),
        ("多頭勝率",
         lambda s: f"{s.get('wr_long', 0):.1f}%"  if "error" not in s else "N/A",
         "wr_long",      "wr"),
        ("空頭勝率",
         lambda s: f"{s.get('wr_short', 0):.1f}%" if "error" not in s else "N/A",
         "wr_short",     "wr"),
        ("累積績效",
         lambda s: f"{s.get('strat_total', 0):+.2f}%" if "error" not in s else "N/A",
         "strat_total",  "pos_red"),
        ("年化報酬",
         lambda s: _ann_text(s),
         "strat_ann",    "pos_red"),
        ("年化波動度",
         lambda s: f"{s.get('strat_vol', 0):.2f}%"     if "error" not in s else "N/A",
         None,           "neutral"),
        ("最大回撤",
         lambda s: f"{s.get('strat_mdd', 0):.2f}%"     if "error" not in s else "N/A",
         "strat_mdd",    "neg_green"),
        ("夏普比率",
         lambda s: f"{s.get('strat_sharpe', 0):.2f}"   if "error" not in s else "N/A",
         "strat_sharpe", "pos_red"),
        ("在市場時間",
         lambda s: f"{s.get('in_mkt', 0):.1f}%"        if "error" not in s else "N/A",
         None, "neutral"),
        ("多單時間",
         lambda s: f"{s.get('in_long', 0):.1f}%"       if "error" not in s else "N/A",
         None, "neutral"),
        ("空單時間",
         lambda s: f"{s.get('in_short', 0):.1f}%"      if "error" not in s else "N/A",
         None, "neutral"),
    ]

    n_rows = 1 + len(ROWS)  # 1個欄位標題列 + 資料列
    tbl_shape = slide.shapes.add_table(
        n_rows, 3,
        Inches(0.00), Inches(4.40),
        Inches(5.41), Inches(6.94)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.00)   # 指標名稱欄
    tbl.columns[1].width = Inches(1.705)  # 近三年欄
    tbl.columns[2].width = Inches(1.705)  # YTD 欄

    # 列高設定：標題列稍高，「回測期間」列雙行需更高
    tbl.rows[0].height = int(Inches(0.40))
    tbl.rows[1].height = int(Inches(0.70))  # 回測期間列雙行
    base_h = int(Inches(0.46))
    remaining_h = int(Inches(6.94) - Inches(0.40) - Inches(0.70))
    per_row_h   = remaining_h // (len(ROWS) - 1)
    for i in range(2, n_rows):
        tbl.rows[i].height = per_row_h

    # ── 顏色常數 ─────────────────────────────────────────────
    C        = RGBColor
    BG_HDR   = (30, 58, 95)
    BG_ALT1  = (241, 245, 249)
    BG_ALT2  = (255, 255, 255)
    FG_WHITE = (255, 255, 255)
    FG_BLUE  = (30, 58, 95)
    FG_RED   = (220, 38, 38)
    FG_GREEN = (22, 163, 74)
    FG_GOLD  = (255, 215, 60)
    FG_GREY  = (100, 116, 139)

    def _w(row, col, text, fg=FG_WHITE, bg=BG_HDR,
           size=14, bold=False, align=PP_ALIGN.CENTER, wrap=False):
        cell = tbl.cell(row, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C(*bg)
        tf = cell.text_frame
        tf.word_wrap = wrap
        tf.margin_top    = Pt(3.6)
        tf.margin_bottom = Pt(3.6)
        tf.margin_left   = Pt(7.2)
        tf.margin_right  = Pt(7.2)
        # 支援「\n」換行（用於回測期間日期）
        lines = text.split('\n')
        for li, line in enumerate(lines):
            if li == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size  = Pt(size)
            run.font.bold  = bold
            run.font.color.rgb = C(*fg)
            run.font.name  = "微軟正黑體"

    def _fg(stats, key, ctype):
        if key is None or "error" in stats:
            return FG_BLUE
        val = stats.get(key, None)
        if val is None:
            return FG_GREY   # 「—（未滿一年）」用灰色
        if ctype == "wr":
            return FG_RED if val >= 60 else FG_BLUE
        elif ctype == "pos_red":
            return FG_RED if val > 0 else (FG_GREEN if val < 0 else FG_BLUE)
        elif ctype == "neg_green":
            return FG_GREEN if val < 0 else FG_BLUE
        return FG_BLUE

    # ── 欄位標題列 (row 0) ───────────────────────────────────
    _w(0, 0, "績效指標",       fg=FG_WHITE, bg=BG_HDR, size=14, bold=True)
    _w(0, 1, "成立至今",       fg=FG_WHITE, bg=BG_HDR, size=14, bold=True)
    _w(0, 2, "今年以來 (YTD)", fg=FG_GOLD,  bg=BG_HDR, size=14, bold=True)

    # ── 資料列 (row 1 起) ────────────────────────────────────
    for r, (label, val_fn, key, ctype) in enumerate(ROWS, start=1):
        bg      = BG_ALT1 if r % 2 == 1 else BG_ALT2
        is_bold = (ctype != "neutral")
        v3  = val_fn(stats_3yr)
        vy  = val_fn(stats_ytd)
        fg3 = _fg(stats_3yr, key, ctype)
        fgy = _fg(stats_ytd, key, ctype)

        # 統一字體大小為 14pt，除了特例縮小
        sz_l = 14
        sz_3 = 14
        sz_y = 14

        # 特例：未滿一年文字縮小
        if v3.startswith("—"): sz_3 = 10
        if vy.startswith("—"): sz_y = 10

        _w(r, 0, label, fg=(55, 65, 81), bg=bg, size=sz_l, bold=False, align=PP_ALIGN.CENTER)
        _w(r, 1, v3, fg=fg3, bg=bg, size=sz_3, bold=is_bold and not v3.startswith("—"))
        _w(r, 2, vy, fg=fgy, bg=bg, size=sz_y, bold=is_bold and not vy.startswith("—"))



def add_factor_table_to_slide(slide, d):
    """
    在 PPT slide 上建立可編輯的因子狀況表（取代原本的 matplotlib 圖片）。
    """
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    i = len(d) - 1;  p = max(0, i - 1)
    def g(c, idx=i):
        try:    return float(d[c].iloc[idx])
        except: return 0.0

    # ── 計算各因子最新數值 ────────────────────────────────────
    twii = g('TWII');      tp   = g('TWII', p)
    twii_chg = (twii/tp-1)*100 if tp > 0 else 0
    bias = g('乖離');      rsi  = g('RSI');  rp = g('RSI', p)
    mom5 = g('MOM5')
    ef   = g('EF');        efp  = g('EF', p); ef20 = g('EF_MA20'); ef60 = g('EF_MA60')
    sox  = g('SOX');       soxp = g('SOX', p)
    sox_chg = (sox/soxp-1)*100 if soxp > 0 else 0
    sox20 = g('SOX_MA20'); sox60 = g('SOX_MA60')
    tsmc = g('TSMC_TW');   tsmcp = g('TSMC_TW', p)
    tsmc_chg = (tsmc/tsmcp-1)*100 if tsmcp > 0 else 0
    tvol = g('TSMC_Vol');  tvma  = g('TS_VolMA')
    vr   = tvol/tvma if tvma > 0 else 0
    adrz = g('ADR_Z')
    fi   = g('FI_Net');    fi5 = g('FI_5MA');  fiz = g('FI_Z')

    # ── 表格內容 ─────────────────────────────────────────────
    # ('HDR',) 表欄位標題列; ('GROUP', 標題) 表分組列; ('DATA', 因子, 值, 狀態, 顏色)
    TABLE_DATA = [
        ('HDR',),
        ('GROUP', '加權指數 (TAIEX)'),
        ('DATA', '收盤指數',   f"{twii:,.0f}",   f"{twii_chg:+.2f}%",
         'red' if twii_chg > 0 else 'green'),
        ('DATA', '季線乖離率', f"{bias:+.2f}%",
         '超賣' if bias < -8 else ('超買' if bias > 8 else '常態'),
         'red' if bias < -8 else ('green' if bias > 8 else 'neutral')),
        ('DATA', 'RSI(14)',    f"{rsi:.1f}",      f"{rsi - rp:+.1f}",
         'red' if rsi < 40 else ('green' if rsi > 55 else 'neutral')),
        ('DATA', '5日動量',    f"{mom5:+.2f}%",
         '強多' if mom5 > 2 else ('強空' if mom5 < -2 else '中性'),
         'red' if mom5 > 2 else ('green' if mom5 < -2 else 'neutral')),
        ('GROUP', '風格輪動 (電金比)'),
        ('DATA', '電金比',   f"{ef:.4f}",  f"{ef - efp:+.4f}",
         'red' if ef > efp else 'green'),
        ('DATA', 'vs 月線',  f"{ef20:.4f}", '站上' if ef > ef20 else '跌破',
         'red' if ef > ef20 else 'green'),
        ('DATA', 'vs 季線',  f"{ef60:.4f}", '站上' if ef > ef60 else '跌破',
         'red' if ef > ef60 else 'green'),
        ('GROUP', '半導體 (海外+本土)'),
        ('DATA', '費城半導體', f"{sox:,.0f}",  f"{sox_chg:+.2f}%",
         'red' if sox_chg > 0 else 'green'),
        ('DATA', '台積電現貨', f"{tsmc:,.0f}", f"{tsmc_chg:+.2f}%",
         'red' if tsmc_chg > 0 else 'green'),
        ('DATA', '台積量比',   f"{vr:.2f}x",
         '爆量' if vr > 1.5 else '正常',
         'green' if vr > 1.5 else 'neutral'),
        ('GROUP', '海外資金面'),
        ('DATA', 'ADR Z-score',   f"{adrz:+.2f}σ",
         '搶單' if adrz > 0.8 else ('撤退' if adrz < -1.0 else '中性'),
         'red' if adrz > 0.8 else ('green' if adrz < -1.0 else 'neutral')),
        ('DATA', '外資今日(億)', f"{fi:+.1f}",
         '買超' if fi > 0 else '賣超',
         'red' if fi > 0 else 'green'),
        ('DATA', '外資 5MA(億)', f"{fi5:+.1f}",
         '連買' if fi5 > 0 else '連賣',
         'red' if fi5 > 0 else 'green'),
        ('DATA', '外資 Z-score', f"{fiz:+.2f}σ",
         '強買超' if fiz > 1.2 else ('強賣超' if fiz < -1.2 else '中性'),
         'red' if fiz > 1.2 else ('green' if fiz < -1.2 else 'neutral')),
    ]

    n_rows = len(TABLE_DATA)
    tbl_shape = slide.shapes.add_table(
        n_rows, 3,
        Inches(5.70), Inches(4.40),
        Inches(5.30), Inches(6.95)
    )
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.00)
    tbl.columns[1].width = Inches(1.65)
    tbl.columns[2].width = Inches(1.65)

    # 列高設定
    hdr_h  = int(Inches(0.349))
    grp_h  = int(Inches(0.489))
    data_h = int(Inches(0.332))
    for r, row_data in enumerate(TABLE_DATA):
        if row_data[0] == 'HDR':    tbl.rows[r].height = hdr_h
        elif row_data[0] == 'GROUP': tbl.rows[r].height = grp_h
        else:                        tbl.rows[r].height = data_h

    # ── 顏色常數 ─────────────────────────────────────────────
    C        = RGBColor
    BG_HDR   = (30, 58, 95)
    BG_GROUP = (51, 65, 85)
    BG_ALT1  = (250, 250, 250)
    BG_ALT2  = (255, 255, 255)
    FG_WHITE = (255, 255, 255)
    FG_BLUE  = (30, 58, 95)
    FG_RED   = (220, 38, 38)
    FG_GREEN = (22, 163, 74)
    FG_GREY  = (100, 116, 139)
    CMAP     = {'red': FG_RED, 'green': FG_GREEN, 'neutral': FG_GREY}

    def _w(row, col, text, fg=FG_WHITE, bg=BG_HDR,
           size=14, bold=False, align=PP_ALIGN.CENTER):
        cell = tbl.cell(row, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C(*bg)
        tf = cell.text_frame
        tf.word_wrap = False
        tf.margin_top    = Pt(3.6)
        tf.margin_bottom = Pt(3.6)
        tf.margin_left   = Pt(7.2)
        tf.margin_right  = Pt(7.2)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = C(*fg)
        run.font.name  = "微軟正黑體"

    data_cnt = 0
    for r, row_data in enumerate(TABLE_DATA):
        rtype = row_data[0]
        if rtype == 'HDR':
            _w(r, 0, '因子',   fg=FG_WHITE, bg=BG_HDR, size=14, bold=True)
            _w(r, 1, '當前值', fg=FG_WHITE, bg=BG_HDR, size=14, bold=True)
            _w(r, 2, '狀態',   fg=FG_WHITE, bg=BG_HDR, size=14, bold=True)
        elif rtype == 'GROUP':
            for col in range(3):
                _w(r, col, row_data[1] if col == 0 else '',
                   fg=FG_WHITE, bg=BG_GROUP, size=14, bold=True,
                   align=PP_ALIGN.CENTER)  # 分組標題也置中
        else:
            _, label, val, status, color = row_data
            bg = BG_ALT1 if data_cnt % 2 == 0 else BG_ALT2
            _w(r, 0, label,  fg=(55, 65, 81), bg=bg, size=14, align=PP_ALIGN.CENTER)
            _w(r, 1, val,    fg=FG_BLUE, bg=bg, size=14, bold=True)
            _w(r, 2, status, fg=CMAP.get(color, FG_GREY), bg=bg, size=14, bold=True)
            data_cnt += 1



def make_chart(d, bt, out_path):
    cutoff  = d["date"].iloc[-1]-pd.Timedelta(days=365)
    sub     = d[d["date"]>=cutoff].copy().reset_index(drop=True)
    exp_sub = bt["exp"][(d["date"]>=cutoff).values]
    fig, ax = plt.subplots(figsize=(14,5.0),facecolor='white')
    ax.set_facecolor('#F9FAFB')
    in_long=False; in_short=False; ls_dt=sub["date"].iloc[0]
    for i in range(len(sub)):
        e=exp_sub[i]; dt=sub["date"].iloc[i]
        if e==1 and not in_long:
            ls_dt=dt; in_long=True; in_short=False
        elif e==-1 and not in_short:
            ls_dt=dt; in_short=True; in_long=False
        elif e==0 and (in_long or in_short):
            # ★ 多單背景=淺紅, 空單背景=淺綠
            ax.axvspan(ls_dt,dt,color='#FEE2E2' if in_long else '#DCFCE7',
                       alpha=0.55,zorder=1,lw=0)
            in_long=False; in_short=False
    if in_long or in_short:
        ax.axvspan(ls_dt,sub["date"].iloc[-1],
                   color='#FEE2E2' if in_long else '#DCFCE7',alpha=0.55,zorder=1,lw=0)
    ax.plot(sub["date"],sub["MA60"],color='#F43F5E',linewidth=1.4,
            linestyle='--',alpha=0.75,zorder=2)
    ax.plot(sub["date"],sub["MA10"],color='#F59E0B',linewidth=1.1,
            linestyle=(0,(3,2)),alpha=0.75,zorder=2)
    ax.plot(sub["date"],sub["TWII"],color='#1E40AF',linewidth=2.4,
            zorder=3,solid_capstyle='round')
    twii_rng=sub["TWII"].max()-sub["TWII"].min(); off=twii_rng*0.016
    trades_1y=[t for t in bt["trades"] if t["entry_date"]>=cutoff]
    for t in trades_1y:
        is_open="exit_date" not in t
        if t["dir_code"]==1:
            # ★ 多單進場：紅色▲
            ax.scatter(t["entry_date"],t["entry_price"]-off,
                       marker='^',s=170,color='#DC2626',edgecolor='white',
                       linewidth=1.8,zorder=6)
            if is_open:
                ax.annotate('持倉中',xy=(t["entry_date"],t["entry_price"]-off),
                            xytext=(10,18),textcoords='offset points',
                            fontsize=9,color='#B91C1C',fontweight='bold',
                            bbox=dict(boxstyle='round,pad=0.2',facecolor='#FEE2E2',
                                      edgecolor='#DC2626',linewidth=1))
            else:
                # ★ 多單出場：淺紅色○
                ax.scatter(t["exit_date"],t["exit_price"],
                           marker='o',s=75,color='#FCA5A5',edgecolor='#DC2626',
                           linewidth=1.5,zorder=6)
        else:
            # ★ 空單進場：綠色▼
            ax.scatter(t["entry_date"],t["entry_price"]+off,
                       marker='v',s=170,color='#16A34A',edgecolor='white',
                       linewidth=1.8,zorder=6)
            if is_open:
                ax.annotate('持倉中',xy=(t["entry_date"],t["entry_price"]+off),
                            xytext=(10,-22),textcoords='offset points',
                            fontsize=9,color='#15803D',fontweight='bold',
                            bbox=dict(boxstyle='round,pad=0.2',facecolor='#DCFCE7',
                                      edgecolor='#16A34A',linewidth=1))
            else:
                # ★ 空單出場：淺綠色○
                ax.scatter(t["exit_date"],t["exit_price"],
                           marker='o',s=75,color='#4ADE80',edgecolor='#16A34A',
                           linewidth=1.5,zorder=6)
    legend_elems=[
        Line2D([0],[0],color='#1E40AF',lw=2.4,label='加權指數'),
        Line2D([0],[0],color='#F43F5E',lw=1.4,ls='--',label='季線(60MA)'),
        Line2D([0],[0],color='#F59E0B',lw=1.1,ls=(0,(3,2)),label='MA10'),
        Line2D([0],[0],marker='^',color='w',markerfacecolor='#DC2626',
               markersize=13,label='多單進場'),
        Line2D([0],[0],marker='o',color='w',markerfacecolor='#FCA5A5',
               markeredgecolor='#DC2626',markersize=10,label='多單出場'),
        Line2D([0],[0],marker='v',color='w',markerfacecolor='#16A34A',
               markersize=13,label='空單進場'),
        Line2D([0],[0],marker='o',color='w',markerfacecolor='#4ADE80',
               markeredgecolor='#16A34A',markersize=10,label='空單出場'),
    ]
    ax.legend(handles=legend_elems,loc='upper left',fontsize=12,ncol=7,
              frameon=True,framealpha=0.92,edgecolor='#D1D5DB',
              bbox_to_anchor=(0.0,1.0),handlelength=1.5,columnspacing=0.7)
    ymin=sub["TWII"].min()*0.98; ymax=sub["TWII"].max()*1.02
    ax.set_ylim(ymin,ymax)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x,_: f"{x:,.0f}"))
    ax.yaxis.set_major_locator(mticker.MaxNLocator(nbins=6,prune='both'))
    ax.xaxis.set_major_formatter(mdates.DateFormatter('%m/%Y'))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    plt.setp(ax.xaxis.get_majorticklabels(),rotation=0,ha='center',fontsize=12)
    ax.tick_params(axis='y',labelsize=12,colors='#374151',length=3)
    ax.tick_params(axis='x',labelsize=12,colors='#374151',length=3)
    ax.set_ylabel('加權指數',fontsize=13,color='#374151',labelpad=8)
    ax.yaxis.grid(True,alpha=0.4,linestyle='--',linewidth=0.6,color='#CBD5E1')
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E5E7EB'); ax.spines['bottom'].set_color('#E5E7EB')
    plt.tight_layout(pad=0.8)
    plt.savefig(out_path,dpi=150,facecolor='white',bbox_inches='tight')
    plt.close()


# ═══════════════════════════════════════════════════════════════
# 5. PPT 套版生成
# ═══════════════════════════════════════════════════════════════
def generate_pptx(run_date, state, stats_3yr, stats_ytd, d, chart_img, output_path):
    """生成 PPTX 日報（績效統計表與因子表均為可編輯的 PPT 原生表格）"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        raise ImportError("請先安裝 python-pptx：pip install python-pptx")
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"模板不存在：{TEMPLATE_PATH}")
    prs = Presentation(str(TEMPLATE_PATH))
    slide = prs.slides[0]

    # ── 更新日期與策略建議文字 ──────────────────────────────
    for shape in slide.shapes:
        if shape.name == '群組 28':
            for sub in shape.shapes:
                if sub.has_text_frame and (
                        str(run_date.year) in sub.text_frame.text or
                        '/' in sub.text_frame.text):
                    paras = list(sub.text_frame.paragraphs)
                    for p_idx, p in enumerate(paras):
                        if not p.runs: continue
                        if   p_idx == 0: p.runs[0].text = str(run_date.year)
                        elif p_idx == 1: p.runs[0].text = f"{run_date.month:02d}/{run_date.day:02d}"
        elif '文字方塊 4' in shape.name:
            for p in shape.text_frame.paragraphs:
                if len(p.runs) >= 7:
                    rec = state["recommendation"]; bias = state["bias"]
                    p.runs[2].text = rec + "　　　　 　　"
                    if   rec in ("多單續抱", "多單建倉", "空單轉多"):
                        p.runs[2].font.color.rgb = RGBColor.from_string("DC2626")
                    elif rec in ("空單續抱", "空單建倉", "多單轉空"):
                        p.runs[2].font.color.rgb = RGBColor.from_string("16A34A")
                    else:
                        p.runs[2].font.color.rgb = RGBColor.from_string("1E293B")
                    p.runs[6].text = bias
                    if   bias == "偏多": p.runs[6].font.color.rgb = RGBColor.from_string("DC2626")
                    elif bias == "偏空": p.runs[6].font.color.rgb = RGBColor.from_string("16A34A")
                    else:                p.runs[6].font.color.rgb = RGBColor.from_string("1E293B")

    _EMU = 914400
    for _s in slide.shapes:
        if _s.name == '文字方塊 29':
            _s.top = int(11.64 * _EMU); _s.height = int(4.61 * _EMU); break

    # ── 折線圖（圖片，維持不變）──────────────────────────────
    slide.shapes.add_picture(chart_img, Inches(-0.08), Inches(12.25),
                             width=Inches(11.08), height=Inches(3.95))

    # ── 可編輯績效統計表（左）────────────────────────────────
    add_stats_table_to_slide(slide, stats_3yr, stats_ytd)

    # ── 可編輯因子狀況表（右）────────────────────────────────
    add_factor_table_to_slide(slide, d)

    # ── 得分說明列 ────────────────────────────────────────────
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(3.50), Inches(10.5), Inches(0.30))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = (f"加權收盤 {state['last_close']:,.0f}  |  "
              f"多頭得分 ml={state['ml']:.2f} (門檻 {LE})  |  "
              f"空頭得分 ms={state['ms']:.2f} (門檻 {SE})  |  "
              f"大環境門票：多 {'OK' if state['gL'] else 'NG'}  "
              f"空 {'OK' if state['gS'] else 'NG'}")
    r.font.size = Pt(11); r.font.name = "Microsoft JhengHei"
    r.font.color.rgb = RGBColor.from_string("475569")

    prs.save(str(output_path))



# ═══════════════════════════════════════════════════════════════
# 5-B. Excel 進出場明細產生
# ═══════════════════════════════════════════════════════════════
def export_trades_excel(d, bt, run_date, output_dir) -> Path:
    """
    將近一年的進出場明細輸出為 Excel 檔案。
    欄位：方向、進場日期、進場點位、出場日期、出場點位、持倉天數、損益(%)、勝/負

    ★ v2.9：不再對交易做「成熟/暫定」篩選或標記——因為 apply_signal_cache()
      已經確保每一天的 ml/ms 分數在第一次被算出來的當下就永久鎖定，
      這裡看到的每一筆交易，就是「那一天實際被算出來、實際會拿去操作」
      的決策結果，本來就不會再變動，不需要也不應該過濾掉任何一筆。
    """
    try:
        import openpyxl
        from openpyxl.styles import (PatternFill, Font, Alignment,
                                     Border, Side, numbers)
        from openpyxl.utils import get_column_letter
    except ImportError:
        print("  ⚠️  openpyxl 未安裝，跳過 Excel 輸出")
        return None

    cutoff = d["date"].iloc[-1] - pd.Timedelta(days=365)
    trades_1y = [t for t in bt["trades"] if t["entry_date"] >= cutoff]

    # ── 整理成 DataFrame ──────────────────────────────────────
    rows = []
    for t in trades_1y:
        is_open  = "exit_date" not in t
        dir_str  = "多單" if t["dir_code"] == 1 else "空單"
        rows.append({
            "方向":       dir_str,
            "進場日期":   t["entry_date"].strftime("%Y/%m/%d"),
            "進場點位":   round(t["entry_price"], 0),
            "出場日期":   t["exit_date"].strftime("%Y/%m/%d") if not is_open else "持倉中",
            "出場點位":   round(t["exit_price"], 0)           if not is_open else "-",
            "持倉天數":   t["n_days"]                         if not is_open else "-",
            "損益(%)":    round(t["pct_return"], 2)            if not is_open else "-",
            "結果":       ("獲利 ✅" if t["is_win"] else "虧損 ❌") if not is_open else "持倉中 🔄",
        })

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "近一年進出場明細"

    # ── 標題列 ───────────────────────────────────────────────
    title = f"台指策略日報 — 近一年進出場明細（截至 {run_date.strftime('%Y/%m/%d')}）"
    ws.merge_cells("A1:H1")
    ws["A1"] = title
    ws["A1"].font      = Font(name="微軟正黑體", size=14, bold=True, color="FFFFFF")
    ws["A1"].fill      = PatternFill("solid", fgColor="1E3A5F")
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 32

    # ── 欄位標題 ─────────────────────────────────────────────
    headers = ["方向", "進場日期", "進場點位", "出場日期", "出場點位",
               "持倉天數", "損益(%)", "結果"]
    thin    = Side(style="thin", color="D1D5DB")
    border  = Border(left=thin, right=thin, top=thin, bottom=thin)
    for col_idx, h in enumerate(headers, start=1):
        cell           = ws.cell(row=2, column=col_idx, value=h)
        cell.font      = Font(name="微軟正黑體", size=11, bold=True, color="FFFFFF")
        cell.fill      = PatternFill("solid", fgColor="334155")
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border    = border
    ws.row_dimensions[2].height = 22

    # ── 資料列 ───────────────────────────────────────────────
    fill_long  = PatternFill("solid", fgColor="FEE2E2")   # 淡紅（多單）
    fill_short = PatternFill("solid", fgColor="DCFCE7")   # 淡綠（空單）
    fill_open  = PatternFill("solid", fgColor="FEF9C3")   # 淡黃（持倉中）
    font_win   = Font(name="微軟正黑體", size=11, color="DC2626", bold=True)
    font_lose  = Font(name="微軟正黑體", size=11, color="16A34A", bold=True)
    font_base  = Font(name="微軟正黑體", size=11)

    for row_idx, row in enumerate(rows, start=3):
        is_holding = (row["出場日期"] == "持倉中")
        is_long    = (row["方向"] == "多單")
        bg_fill    = fill_open if is_holding else (fill_long if is_long else fill_short)
        pnl        = row["損益(%)"]

        for col_idx, key in enumerate(headers, start=1):
            cell           = ws.cell(row=row_idx, column=col_idx, value=row[key])
            cell.fill      = bg_fill
            cell.border    = border
            cell.alignment = Alignment(horizontal="center", vertical="center")
            if key == "損益(%)" and isinstance(pnl, float):
                cell.font = font_win if pnl >= 0 else font_lose
            elif key == "結果":
                if "獲利" in str(row[key]):  cell.font = font_win
                elif "虧損" in str(row[key]): cell.font = font_lose
                else: cell.font = font_base
            else:
                cell.font = font_base
        ws.row_dimensions[row_idx].height = 20

    # ── 統計摘要列 ───────────────────────────────────────────
    done_trades = [r for r in rows if r["出場日期"] != "持倉中"]
    n_total  = len(done_trades)
    n_win    = sum(1 for r in done_trades if "獲利" in str(r["結果"]))
    n_long   = sum(1 for r in done_trades if r["方向"] == "多單")
    n_short  = sum(1 for r in done_trades if r["方向"] == "空單")
    w_long   = sum(1 for r in done_trades if r["方向"] == "多單" and "獲利" in str(r["結果"]))
    w_short  = sum(1 for r in done_trades if r["方向"] == "空單" and "獲利" in str(r["結果"]))
    wr_all   = n_win   / n_total  * 100 if n_total  else 0
    wr_long  = w_long  / n_long   * 100 if n_long   else 0
    wr_short = w_short / n_short  * 100 if n_short  else 0
    total_pnl = sum(r["損益(%)"] for r in done_trades if isinstance(r["損益(%)"], float))

    ws.append([])
    summary_row = ws.max_row + 1
    ws.merge_cells(f"A{summary_row}:H{summary_row}")
    ws[f"A{summary_row}"] = ("　".join([
        f"統計摘要（已出場 {n_total} 筆）",
        f"整體勝率：{wr_all:.1f}%",
        f"多單勝率：{wr_long:.1f}%（{n_long} 筆）",
        f"空單勝率：{wr_short:.1f}%（{n_short} 筆）",
        f"累積損益：{total_pnl:+.2f}%",
    ]))
    ws[f"A{summary_row}"].font      = Font(name="微軟正黑體", size=11, bold=True, color="1E3A5F")
    ws[f"A{summary_row}"].fill      = PatternFill("solid", fgColor="E0F2FE")
    ws[f"A{summary_row}"].alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[summary_row].height = 24

    # ── 欄寬調整 ─────────────────────────────────────────────
    col_widths = [10, 14, 14, 14, 14, 12, 12, 14]
    for i, w in enumerate(col_widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w

    # ── 凍結標題 ─────────────────────────────────────────────
    ws.freeze_panes = "A3"

    # ── 儲存 ────────────────────────────────────────────────
    excel_fname = f"台指策略進出場明細_{run_date.strftime('%Y.%m.%d')}.xlsx"
    excel_path  = output_dir / excel_fname
    wb.save(str(excel_path))
    print(f"  ✅  進出場明細已產出：{excel_path}（共 {len(rows)} 筆）")
    return excel_path


def generate_daily_report():
    # ★ 使用台灣時間取得日期
    run_date = datetime.now(TW_TZ).date()
    print("="*60)
    print(f"  台指策略日報  v15 多因子量化模型  v2.9")
    print(f"  執行日期：{run_date}  ({datetime.now(TW_TZ).strftime('%H:%M:%S')} 台灣時間)")
    print("="*60)

    # ★ v2.8：若指定 REBUILD_CACHE=1 則清空快取重建
    reset_cache_if_requested()

    df=fetch_market_data()
    # ★ v2.8：行情快照凍結（防 yfinance 除權息回溯調整 / FinMind 事後校正）
    df=freeze_market_data(df)

    print(f"✓ 資料：{len(df)} 個交易日  "
          f"({df['date'].iloc[0].date()} ~ {df['date'].iloc[-1].date()})")
    start_str=(df["date"].iloc[0]-timedelta(days=10)).strftime("%Y-%m-%d")
    fi=fetch_foreign_investor(start_str)
    if fi.empty and 'FI_Net_csv' in df.columns:
        fi=df.set_index('date')['FI_Net_csv'].dropna()
        print(f"✓ 外資：CSV fallback（{len(fi)} 筆）")
    elif not fi.empty:
        print(f"✓ 外資：FinMind（{len(fi)} 筆）")

    d=build_factors(df,fi)
    ml,ms,gL,gS=compute_scores(d)
    # ★ v2.8：訊號分數鎖定（第一次算出來的當下即永久鎖定，見 v2.9 說明）
    ml,ms=apply_signal_cache(d,ml,ms)

    bt=backtest(d,ml,ms,gL,gS)
    state=determine_state(d,bt,ml,ms,gL,gS)
    stats=compute_1yr_stats(d,bt)      # 近一年摘要（終端機輸出用）
    perf=compute_perf_stats(d,bt)      # 近三年 + YTD（PPT 表格用）

    i_last=len(d)-1
    print(f"\n  策略建議：{state['recommendation']}　｜　模型分析：{state['bias']}")
    print(f"  加權收盤：{state['last_close']:,.0f}（資料日期：{state['last_date']}）")
    print(f"  月線 MA20：{d['MA20'].iloc[i_last]:,.2f}　"
          f"季線 MA60：{d['MA60'].iloc[i_last]:,.2f}")
    print(f"  ADR 折溢價：{d['ADR_Premium'].iloc[i_last]:+.2f}%")
    print(f"  ml={state['ml']:.2f}（門檻{LE}）　ms={state['ms']:.2f}（門檻{SE}）")
    print(f"  門票　多：{'OK' if state['gL'] else 'NG'}　空：{'OK' if state['gS'] else 'NG'}")
    if "error" not in stats:
        print(f"\n  近一年　{stats['n_trades']} 筆  勝率 {stats['wr_all']:.1f}%  "
              f"累積 {stats['strat_total']:+.2f}%  "
              f"MDD {stats['strat_mdd']:.2f}%  Sharpe {stats['strat_sharpe']:.2f}")
    s3 = perf["3yr"]
    if "error" not in s3:
        print(f"  近三年　{s3['n_trades']} 筆  勝率 {s3['wr_all']:.1f}%  "
              f"累積 {s3['strat_total']:+.2f}%  "
              f"MDD {s3['strat_mdd']:.2f}%  Sharpe {s3['strat_sharpe']:.2f}")

    open_trades=[t for t in bt["trades"] if "exit_date" not in t]
    if open_trades:
        ot=open_trades[-1]
        dir_str="多單" if ot["dir_code"]==1 else "空單"
        print(f"\n  ★ 進行中部位：{dir_str}  進場 {ot['entry_date'].date()}"
              f"  進場價 {ot['entry_price']:,.0f}")

    OUTPUT_DIR.mkdir(exist_ok=True)
    chart_img = str(OUTPUT_DIR/"_tmp_chart.png")   # ★ 只剩折線圖需要暫存
    print("\n  生成折線圖...", end='', flush=True)
    make_chart(d, bt, chart_img)
    print(" 完成")

    fname=f"台股策略日報_{run_date.strftime('%Y.%m.%d')}.pptx"
    output_path=OUTPUT_DIR/fname
    print(f"  套版輸出...", end='', flush=True)
    generate_pptx(run_date, state, perf["3yr"], perf["ytd"], d, chart_img, output_path)
    print(" 完成")

    for _p in (chart_img,):
        try: Path(_p).unlink()
        except: pass

    # ── 產生 Excel 進出場明細 ──────────────────────────────────
    print(f"\n  產生進出場明細 Excel...", end='', flush=True)
    excel_path = export_trades_excel(d, bt, run_date, OUTPUT_DIR)
    print(" 完成" if excel_path else " 跳過（openpyxl 未安裝）")

    print(f"\n{'='*60}")
    print(f"  ✅  日報已產出：{output_path}")
    if excel_path:
        print(f"  ✅  進出場明細：{excel_path}")
    print(f"{'='*60}")
    # ★ 回傳 (PPTX路徑, Excel路徑, 策略狀態)，供 run_and_send.py 組 Email 內文
    return output_path, excel_path, state


if __name__=="__main__" or "ipykernel" in sys.modules:
    generate_daily_report()
